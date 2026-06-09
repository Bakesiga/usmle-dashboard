"""
Build CVS Physiology Review deck (80 question/answer slides + title + divider).
Generates matplotlib figures for image-bearing items, then assembles a 16:9 pptx
in a clean USMLE-mimic visual style.

Rules: no em dashes anywhere, original-prose only (from JSONL), 16:9 layout.
"""
import json
import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from pptx.dml.color import RGBColor

BASE = Path("/Users/allanbakesiga/.claude/usmle-dashboard/outputs/study-notes")
JSONL = BASE / "cvs-physiology-uworld-bank.jsonl"
FIG_DIR = BASE / "cvs-physiology-deck-figures"
OUT_PPTX = BASE / "cvs-physiology-review-deck.pptx"

FIG_DIR.mkdir(parents=True, exist_ok=True)

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x10, 0x10, 0x10)
GRAY_LINE = RGBColor(0xC8, 0xCD, 0xD3)
GRAY_LIGHT = RGBColor(0xF3, 0xF5, 0xF7)
GRAY_MID = RGBColor(0x8A, 0x92, 0x9A)
GRAY_TEXT = RGBColor(0x33, 0x3A, 0x40)
USMLE_RED = RGBColor(0xC4, 0x2B, 0x2B)
USMLE_GREEN = RGBColor(0x2E, 0x8B, 0x57)
USMLE_BLUE = RGBColor(0x1F, 0x4E, 0x79)

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"

# Slide dims 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- Figure generation ----------

def _setup_axes(ax):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(direction="out", length=4, colors="#222")
    for s in ("left", "bottom"):
        ax.spines[s].set_color("#222")
    ax.set_facecolor("white")


def _save(fig, path):
    fig.tight_layout()
    fig.savefig(path, dpi=120, facecolor="white", bbox_inches="tight")
    plt.close(fig)


def _pv_loop_points(ax, edv=120, esv=50, sp=120, dp=80, mp=8):
    """Plot a standard PV loop, return dict of corner coords."""
    # Corners
    A = (edv, mp)        # mitral closure (end diastole, bottom-right)
    B = (edv, dp)        # aortic opening (top-right)
    C = (esv, sp)        # aortic closure (top-left)
    D = (esv, mp)        # mitral opening (bottom-left)
    # Isovolumetric contraction A->B (vertical), ejection B->C (curve),
    # isovolumetric relaxation C->D, filling D->A (curve)
    # Ejection: small curve peaking at E ~ ((edv+esv)/2, sp+5)
    Ex = (edv + esv) / 2
    Ey = sp + 5
    # Filling: curve dipping at F ~ ((edv+esv)/2, mp-2)
    Fx = (edv + esv) / 2
    Fy = max(mp - 2, 2)

    # Iso contraction
    ax.plot([A[0], B[0]], [A[1], B[1]], color="black", lw=1.6)
    # Ejection (quadratic via three points)
    t = np.linspace(0, 1, 30)
    x_ej = (1 - t) ** 2 * B[0] + 2 * (1 - t) * t * Ex + t ** 2 * C[0]
    y_ej = (1 - t) ** 2 * B[1] + 2 * (1 - t) * t * Ey + t ** 2 * C[1]
    ax.plot(x_ej, y_ej, color="black", lw=1.6)
    # Iso relaxation
    ax.plot([C[0], D[0]], [C[1], D[1]], color="black", lw=1.6)
    # Filling
    x_fi = (1 - t) ** 2 * D[0] + 2 * (1 - t) * t * Fx + t ** 2 * A[0]
    y_fi = (1 - t) ** 2 * D[1] + 2 * (1 - t) * t * Fy + t ** 2 * A[1]
    ax.plot(x_fi, y_fi, color="black", lw=1.6)

    ax.set_xlabel("LV volume (mL)")
    ax.set_ylabel("LV pressure (mm Hg)")
    _setup_axes(ax)
    return {"A": A, "B": B, "C": C, "D": D, "E": (Ex, Ey), "F": (Fx, Fy)}


def fig_q01(path):
    """Standard PV loop with 6 labeled points A-F per description."""
    fig, ax = plt.subplots(figsize=(5, 3.75))
    pts = _pv_loop_points(ax)
    # Description: E top, C upper-left, B upper-right, A lower-right, D lower-left, F bottom-middle
    label_offsets = {
        "A": (6, -3),
        "B": (6, 3),
        "C": (-8, 5),
        "D": (-9, -3),
        "E": (0, 6),
        "F": (0, -6),
    }
    for k, (x, y) in pts.items():
        dx, dy = label_offsets[k]
        ax.plot(x, y, "o", color="black", markersize=5)
        ax.annotate(k, (x, y), xytext=(x + dx, y + dy), fontsize=11, fontweight="bold",
                    ha="center", va="center")
    ax.set_xlim(30, 160)
    ax.set_ylim(-5, 145)
    ax.set_title("LV pressure volume loop", fontsize=11)
    _save(fig, path)


def _ecg_pqrst(ax, x0, with_p=True, qrs_width=0.04, qrs_amp=1.0, t_amp=0.25,
               invert_t=False, wide=False):
    """Draw one PQRST complex starting at x0. Returns next x position."""
    x = x0
    seg = []
    # baseline
    seg.append((x, 0)); x += 0.05
    seg.append((x, 0))
    if with_p:
        # P wave (small hump)
        for tt in np.linspace(0, 0.08, 8):
            seg.append((x + tt, 0.12 * np.sin(np.pi * tt / 0.08)))
        x += 0.08
    seg.append((x, 0)); x += 0.04
    if wide:
        # Wide QRS
        seg.append((x, 0))
        x += 0.02
        seg.append((x, -0.15 * qrs_amp))
        x += 0.05
        seg.append((x, 0.9 * qrs_amp))
        x += 0.05
        seg.append((x, -0.2 * qrs_amp))
        x += 0.04
        seg.append((x, 0))
    else:
        # Q
        seg.append((x, -0.1 * qrs_amp)); x += 0.01
        # R spike
        seg.append((x, qrs_amp)); x += 0.02
        # S
        seg.append((x, -0.25 * qrs_amp)); x += 0.01
        seg.append((x, 0))
    x += 0.04
    # T wave
    sign = -1 if invert_t else 1
    for tt in np.linspace(0, 0.12, 10):
        seg.append((x + tt, sign * t_amp * np.sin(np.pi * tt / 0.12)))
    x += 0.12
    seg.append((x, 0)); x += 0.05
    xs, ys = zip(*seg)
    ax.plot(xs, ys, color="black", lw=1.2)
    return x


def fig_q02(path):
    """AF strip: irregular narrow QRS, no P waves, chaotic baseline."""
    fig, ax = plt.subplots(figsize=(5.5, 3))
    # Chaotic fibrillatory baseline
    xs_base = np.linspace(0, 5, 1500)
    rng = np.random.default_rng(7)
    base = 0.04 * np.sin(xs_base * 30) + 0.06 * rng.normal(size=len(xs_base)) * 0.5
    ax.plot(xs_base, base, color="black", lw=0.6)
    # Irregular narrow QRS at varying times
    rr = [0.55, 0.78, 0.45, 1.05, 0.62, 0.92, 0.50]
    x = 0.3
    for w in rr:
        if x > 4.7:
            break
        # narrow QRS spike
        ax.plot([x, x + 0.015, x + 0.04, x + 0.06], [0, 1.0, -0.25, 0], color="black", lw=1.4)
        # T
        for tt in np.linspace(0, 0.14, 12):
            ax.plot([x + 0.1 + tt - 0.01, x + 0.1 + tt],
                    [0.25 * np.sin(np.pi * (tt - 0.01) / 0.14),
                     0.25 * np.sin(np.pi * tt / 0.14)], color="black", lw=1.0)
        x += w
    ax.set_xlim(0, 5)
    ax.set_ylim(-0.6, 1.3)
    ax.set_title("Atrial fibrillation (irregularly irregular, no P waves)", fontsize=10)
    ax.set_xlabel("Time (s)")
    ax.set_yticks([])
    _setup_axes(ax)
    _save(fig, path)


def fig_q04(path):
    """ECG strip with three beats X, Y, Z. Y is wide-complex PVC with discordant T, then compensatory pause."""
    fig, ax = plt.subplots(figsize=(5.5, 3))
    x = 0.2
    x = _ecg_pqrst(ax, x); x_X_end = x
    ax.annotate("X", (0.55, 1.15), fontsize=11, fontweight="bold")
    # Y: PVC wide, comes early, T discordant
    x_pvc = x_X_end - 0.05
    x = _ecg_pqrst(ax, x_pvc, with_p=False, wide=True, invert_t=True, qrs_amp=1.1, t_amp=0.35)
    ax.annotate("Y", (x_pvc + 0.15, 1.1), fontsize=11, fontweight="bold")
    # Compensatory pause
    x += 0.45
    # Z normal
    x = _ecg_pqrst(ax, x)
    ax.annotate("Z", (x - 0.35, 1.15), fontsize=11, fontweight="bold")
    ax.set_xlim(0, max(x + 0.2, 3.5))
    ax.set_ylim(-0.6, 1.4)
    ax.set_yticks([])
    ax.set_xlabel("Time (s)")
    ax.set_title("Sinus, PVC, compensatory pause, sinus", fontsize=10)
    _setup_axes(ax)
    _save(fig, path)


def fig_q05(path):
    """Four PV loop panels (A, B, C, D) comparing baseline vs post AV fistula intervention.
    Option B is the correct one per description."""
    fig, axes = plt.subplots(1, 4, figsize=(8.5, 3))
    panels = ["A", "B", "C", "D"]
    for i, ax in enumerate(axes):
        # baseline solid
        _pv_loop_points(ax)
        # Reset ax decoration
        ax.set_xlabel("Vol", fontsize=8)
        ax.set_ylabel("Pres", fontsize=8) if i == 0 else ax.set_ylabel("")
        ax.set_xlim(20, 170)
        ax.set_ylim(-5, 150)
        ax.set_title(panels[i], fontsize=11, fontweight="bold")
        # dashed variant per panel (illustrative shapes)
        if i == 0:  # A: bigger loop higher pressure
            _draw_dashed_loop(ax, edv=130, esv=55, sp=140, dp=85)
        elif i == 1:  # B: correct AV fistula loop, wider, higher EDV, lower ESV, slightly lower SP
            _draw_dashed_loop(ax, edv=140, esv=40, sp=110, dp=72)
        elif i == 2:  # C: narrower
            _draw_dashed_loop(ax, edv=100, esv=55, sp=125, dp=82)
        else:  # D: shifted left lower
            _draw_dashed_loop(ax, edv=110, esv=60, sp=100, dp=70)
        ax.tick_params(labelsize=7)
    fig.suptitle("Pressure volume loops, baseline (solid) vs post (dashed)", fontsize=10)
    _save(fig, path)


def _draw_dashed_loop(ax, edv, esv, sp, dp, mp=8):
    A = (edv, mp); B = (edv, dp); C = (esv, sp); D = (esv, mp)
    Ex = (edv + esv) / 2; Ey = sp + 5
    Fx = (edv + esv) / 2; Fy = max(mp - 2, 2)
    ax.plot([A[0], B[0]], [A[1], B[1]], color="black", lw=1.3, ls=(0, (4, 3)))
    t = np.linspace(0, 1, 30)
    x_ej = (1 - t) ** 2 * B[0] + 2 * (1 - t) * t * Ex + t ** 2 * C[0]
    y_ej = (1 - t) ** 2 * B[1] + 2 * (1 - t) * t * Ey + t ** 2 * C[1]
    ax.plot(x_ej, y_ej, color="black", lw=1.3, ls=(0, (4, 3)))
    ax.plot([C[0], D[0]], [C[1], D[1]], color="black", lw=1.3, ls=(0, (4, 3)))
    x_fi = (1 - t) ** 2 * D[0] + 2 * (1 - t) * t * Fx + t ** 2 * A[0]
    y_fi = (1 - t) ** 2 * D[1] + 2 * (1 - t) * t * Fy + t ** 2 * A[1]
    ax.plot(x_fi, y_fi, color="black", lw=1.3, ls=(0, (4, 3)))


def fig_q07(path):
    """Aortocavitary fistula schematic with chamber boxes and pressure labels."""
    fig, ax = plt.subplots(figsize=(5.5, 4))
    ax.set_xlim(0, 10); ax.set_ylim(0, 8)
    ax.axis("off")
    # Chamber boxes
    chambers = {
        "RA":  (1.0, 4.5, 1.8, 1.5, "RA\n< 5"),
        "RV":  (1.0, 1.5, 1.8, 2.5, "RV\n25/5"),
        "LA":  (6.8, 4.5, 1.8, 1.5, "LA\n< 10"),
        "LV":  (6.8, 1.5, 1.8, 2.5, "LV\n120/10"),
        "Ao":  (4.0, 6.0, 1.8, 1.3, "Aorta\n120/80"),
        "PA":  (3.0, 4.5, 1.4, 1.2, "PA\n25/10"),
    }
    for k, (x, y, w, h, lbl) in chambers.items():
        ax.add_patch(Rectangle((x, y), w, h, fill=False, edgecolor="black", lw=1.2))
        ax.text(x + w / 2, y + h / 2, lbl, ha="center", va="center", fontsize=9)
    # Fistula from Aorta to RV
    ax.annotate("", xy=(2.5, 3.0), xytext=(4.6, 6.0),
                arrowprops=dict(arrowstyle="->", color="#C42B2B", lw=2))
    ax.text(3.5, 4.7, "fistula", color="#C42B2B", fontsize=9, fontweight="bold")
    ax.set_title("Aorta to RV communication", fontsize=10)
    _save(fig, path)


def fig_q08(path):
    """Mitral inflow Doppler-like trace with E and A waves; arrow at deceleration after E."""
    fig, ax = plt.subplots(figsize=(5.5, 3))
    t = np.linspace(0, 1, 500)
    # E wave: peak ~0.18 height 1.0, decel by 0.32
    E = np.exp(-((t - 0.18) ** 2) / (2 * 0.05 ** 2))
    # Diastasis low
    # A wave: peak ~0.78 height 0.55
    A = 0.55 * np.exp(-((t - 0.78) ** 2) / (2 * 0.05 ** 2))
    y = E + A
    ax.plot(t, y, color="black", lw=1.4)
    ax.axhline(0, color="black", lw=0.6)
    ax.annotate("E", (0.18, 1.05), fontsize=11, fontweight="bold", ha="center")
    ax.annotate("A", (0.78, 0.6), fontsize=11, fontweight="bold", ha="center")
    # Arrow on deceleration shoulder
    ax.annotate("", xy=(0.32, 0.18), xytext=(0.42, 0.55),
                arrowprops=dict(arrowstyle="->", color="#C42B2B", lw=1.6))
    ax.text(0.45, 0.58, "decel", color="#C42B2B", fontsize=9)
    ax.set_xlim(0, 1); ax.set_ylim(-0.1, 1.25)
    ax.set_xlabel("Diastole (time)")
    ax.set_ylabel("Mitral inflow")
    ax.set_title("Diastolic mitral inflow, E and A waves", fontsize=10)
    _setup_axes(ax)
    _save(fig, path)


def fig_q11(path):
    """LV pressure and volume vs time with five labeled points."""
    fig, ax1 = plt.subplots(figsize=(5.5, 3.4))
    t = np.linspace(0, 1, 1000)
    # LV pressure: rises sharply (iso contraction) ~0.05, peaks ~0.25, falls ~0.4
    p = np.piecewise(
        t,
        [t < 0.04, (t >= 0.04) & (t < 0.1), (t >= 0.1) & (t < 0.35),
         (t >= 0.35) & (t < 0.45), t >= 0.45],
        [
            lambda x: 8 + 4 * np.sin(np.pi * x / 0.04),  # atrial kick slight rise
            lambda x: 8 + 120 * ((x - 0.04) / 0.06),     # iso contraction
            lambda x: 128 - 30 * np.sin(np.pi * (x - 0.1) / 0.25 - 0.5),  # ejection curve
            lambda x: 100 - 90 * ((x - 0.35) / 0.10),    # iso relaxation
            lambda x: 8 + 4 * np.exp(-(x - 0.5) ** 2 / 0.02),
        ],
    )
    ax1.plot(t, p, color="black", lw=1.5, label="LV pressure")
    ax1.set_ylabel("LV pressure (mm Hg)")
    ax1.set_xlabel("Time (one cardiac cycle)")

    # LV volume on second y axis: 50 ESV, 120 EDV
    ax2 = ax1.twinx()
    v = np.piecewise(
        t,
        [t < 0.04, (t >= 0.04) & (t < 0.1), (t >= 0.1) & (t < 0.35),
         (t >= 0.35) & (t < 0.45), t >= 0.45],
        [
            lambda x: 115 + 5 * (x / 0.04),
            lambda x: 120 * np.ones_like(x),
            lambda x: 120 - 70 * ((x - 0.1) / 0.25),
            lambda x: 50 * np.ones_like(x),
            lambda x: 50 + 65 * (1 - np.exp(-(x - 0.45) / 0.15)),
        ],
    )
    ax2.plot(t, v, color="gray", lw=1.4, ls=(0, (4, 2)), label="LV volume")
    ax2.set_ylabel("LV volume (mL)")

    # Points
    pts = {
        "A": (0.04, p[int(0.04 * 1000)]),
        "B": (0.10, p[int(0.10 * 1000)]),
        "C": (0.35, p[int(0.35 * 1000)]),
        "D": (0.45, p[int(0.45 * 1000)]),
        "E": (0.70, p[int(0.70 * 1000)]),
    }
    for k, (x, y) in pts.items():
        ax1.plot(x, y, "o", color="#C42B2B", markersize=5)
        ax1.annotate(k, (x, y), xytext=(x, y + 8), fontsize=10, fontweight="bold", ha="center")

    ax1.set_xlim(0, 1)
    ax1.set_ylim(-5, 160)
    ax2.set_ylim(20, 160)
    _setup_axes(ax1)
    ax1.set_title("LV pressure and volume vs time", fontsize=10)
    _save(fig, path)


def _slow_ap(ax, x0=0, label=None):
    t = np.linspace(0, 1, 500)
    # Phase 4 slow upslope from -60 to -40, then slow phase 0 to +20, plateau-less, phase 3 back to -60
    v = np.piecewise(
        t,
        [t < 0.35, (t >= 0.35) & (t < 0.45), (t >= 0.45) & (t < 0.7), t >= 0.7],
        [
            lambda x: -60 + 20 * (x / 0.35),                # phase 4
            lambda x: -40 + 60 * ((x - 0.35) / 0.10),       # phase 0 slow
            lambda x: 20 - 80 * ((x - 0.45) / 0.25),        # phase 3
            lambda x: -60 + 0 * x,
        ],
    )
    ax.plot(t + x0, v, color="black", lw=1.6)
    ax.axhline(-60, color="gray", lw=0.6, ls=(0, (2, 2)))
    if label:
        ax.set_title(label, fontsize=10)


def _fast_ap(ax, x0=0, label=None):
    t = np.linspace(0, 1, 500)
    v = np.piecewise(
        t,
        [t < 0.05, (t >= 0.05) & (t < 0.08), (t >= 0.08) & (t < 0.55),
         (t >= 0.55) & (t < 0.75), t >= 0.75],
        [
            lambda x: -90 + 0 * x,
            lambda x: -90 + 120 * ((x - 0.05) / 0.03),   # phase 0 fast
            lambda x: 30 - 15 * ((x - 0.08) / 0.47),     # plateau
            lambda x: 15 - 105 * ((x - 0.55) / 0.20),    # phase 3
            lambda x: -90 + 0 * x,
        ],
    )
    ax.plot(t + x0, v, color="black", lw=1.6)
    ax.axhline(-90, color="gray", lw=0.6, ls=(0, (2, 2)))
    if label:
        ax.set_title(label, fontsize=10)


def fig_q12(path):
    """Pacemaker (slow response) action potential."""
    fig, ax = plt.subplots(figsize=(5, 3))
    _slow_ap(ax, label="Slow response (pacemaker) AP")
    ax.set_xlim(0, 1.05)
    ax.set_ylim(-75, 40)
    ax.set_xlabel("Time")
    ax.set_ylabel("Membrane potential (mV)")
    ax.text(0.05, -55, "phase 4", fontsize=8, color="#555")
    ax.text(0.36, 0, "phase 0", fontsize=8, color="#555")
    ax.text(0.55, -10, "phase 3", fontsize=8, color="#555")
    _setup_axes(ax)
    _save(fig, path)


def fig_q16(path):
    """Phonocardiogram: fixed split S2 across expiration and inspiration (ASD)."""
    fig, axes = plt.subplots(1, 2, figsize=(7, 3), sharey=True)
    for i, (ax, label) in enumerate(zip(axes, ["Expiration", "Inspiration"])):
        t = np.linspace(0, 1, 1200)
        y = np.zeros_like(t)
        # S1 at 0.1
        y += _heartburst(t, 0.10, 0.018, 0.9)
        # Systolic ejection murmur diamond from 0.15 to 0.40, peak ~0.22
        ej = np.where((t > 0.15) & (t < 0.40),
                      0.35 * np.sin(np.pi * (t - 0.15) / 0.25) * np.sin(80 * np.pi * (t - 0.15)),
                      0)
        y += ej
        # A2 at 0.5, P2 at 0.58 (fixed split, no change between panels)
        y += _heartburst(t, 0.50, 0.012, 0.7)
        y += _heartburst(t, 0.58, 0.012, 0.55)
        ax.plot(t, y, color="black", lw=0.9)
        ax.set_title(label, fontsize=10)
        ax.set_xlim(0, 1); ax.set_ylim(-1.0, 1.0)
        ax.set_xlabel("Time")
        ax.text(0.10, 0.95, "S1", fontsize=8, ha="center")
        ax.text(0.22, 0.85, "SM", fontsize=8, ha="center")
        ax.text(0.50, 0.85, "A2", fontsize=8, ha="center")
        ax.text(0.58, 0.75, "P2", fontsize=8, ha="center")
        _setup_axes(ax)
        ax.set_yticks([])
    fig.suptitle("Fixed split S2 across respiration", fontsize=10)
    _save(fig, path)


def _heartburst(t, center, width, amp):
    env = np.exp(-((t - center) ** 2) / (2 * width ** 2))
    osc = np.sin(2 * np.pi * 90 * (t - center))
    return amp * env * osc


def fig_q21(path):
    """Cardiac function curve and venous return curve crossing at steady state."""
    fig, ax = plt.subplots(figsize=(5, 3.5))
    rap = np.linspace(-2, 12, 200)
    cf = 7 * (1 - np.exp(-(rap + 2) / 4))  # Frank Starling rising
    vr = 5 - 0.4 * rap                     # venous return descending
    vr = np.where(vr < 0, 0, vr)
    ax.plot(rap, cf, color="#C42B2B", lw=1.8, label="Cardiac function")
    ax.plot(rap, vr, color="#1F4E79", lw=1.8, label="Venous return")
    # Steady state intersection
    diffs = np.abs(cf - vr)
    idx = np.argmin(diffs)
    ax.plot(rap[idx], cf[idx], "o", color="black", markersize=6)
    ax.annotate("steady state", (rap[idx], cf[idx]),
                xytext=(rap[idx] + 1.5, cf[idx] + 0.5), fontsize=9,
                arrowprops=dict(arrowstyle="->", lw=0.8))
    ax.set_xlabel("Right atrial pressure (mm Hg)")
    ax.set_ylabel("Cardiac output / venous return (L/min)")
    ax.legend(fontsize=8, frameon=False)
    ax.set_xlim(-2, 12); ax.set_ylim(0, 8)
    _setup_axes(ax)
    ax.set_title("Cardiac function and venous return curves", fontsize=10)
    _save(fig, path)


def fig_q24(path):
    """Two AP tracings stacked: slow (top), fast (bottom)."""
    fig, axes = plt.subplots(2, 1, figsize=(5, 4.2))
    _slow_ap(axes[0], label="Slow response (pacemaker)")
    axes[0].set_ylim(-75, 40); axes[0].set_xlim(0, 1.05)
    axes[0].set_ylabel("mV")
    _setup_axes(axes[0])
    _fast_ap(axes[1], label="Fast response (ventricular myocyte)")
    axes[1].set_ylim(-100, 50); axes[1].set_xlim(0, 1.05)
    axes[1].set_ylabel("mV"); axes[1].set_xlabel("Time")
    _setup_axes(axes[1])
    _save(fig, path)


def fig_q34(path):
    """Autonomic receptor signaling reference table."""
    fig, ax = plt.subplots(figsize=(6.2, 3.5))
    ax.axis("off")
    data = [
        ["alpha 1", "IP3 up", "vasoconstriction"],
        ["alpha 2", "cAMP down", "CNS sympatholytic"],
        ["beta 1",  "cAMP up", "cardiac inotropy, chronotropy"],
        ["beta 2",  "cAMP up", "vasodilation, bronchodilation"],
        ["M1",      "IP3 up", "intestinal motility"],
        ["M2",      "cAMP down", "decreased HR, contractility"],
        ["M3",      "IP3 up", "bronchoconstriction, glands"],
    ]
    columns = ["Receptor", "Second msgr", "Effects"]
    table = ax.table(cellText=data, colLabels=columns, loc="center", cellLoc="left",
                     colWidths=[0.20, 0.27, 0.53],
                     colColours=["#F0F0F0"] * 3)
    table.auto_set_font_size(False)
    table.set_fontsize(8.5)
    table.scale(1, 1.35)
    ax.set_title("Autonomic receptor signaling reference", fontsize=10)
    _save(fig, path)


def fig_q35(path):
    """Phonocardiogram: holosystolic murmur louder on inspiration (TR)."""
    fig, axes = plt.subplots(1, 2, figsize=(7, 3), sharey=True)
    for i, (ax, label, amp_factor) in enumerate(zip(axes, ["Expiration", "Inspiration"], [0.4, 0.75])):
        t = np.linspace(0, 1, 1200)
        y = np.zeros_like(t)
        y += _heartburst(t, 0.10, 0.018, 0.85)  # S1
        # Holosystolic murmur from 0.13 to 0.50
        mu = np.where((t > 0.13) & (t < 0.50),
                      amp_factor * np.sin(80 * np.pi * (t - 0.13)) *
                      np.exp(-((t - 0.30) / 0.25) ** 2),
                      0)
        y += mu
        y += _heartburst(t, 0.52, 0.013, 0.65)  # S2
        ax.plot(t, y, color="black", lw=0.9)
        ax.set_title(label, fontsize=10)
        ax.set_xlim(0, 1); ax.set_ylim(-1.0, 1.0)
        ax.set_xlabel("Time"); ax.set_yticks([])
        ax.text(0.10, 0.92, "S1", fontsize=8, ha="center")
        ax.text(0.30, 0.85, "holosystolic", fontsize=8, ha="center")
        ax.text(0.52, 0.85, "S2", fontsize=8, ha="center")
        _setup_axes(ax)
    fig.suptitle("Murmur intensifies with inspiration", fontsize=10)
    _save(fig, path)


def fig_q36(path):
    """Coronary autoregulation: flow vs perfusion pressure with autoreg plateau, baseline and four post-stenosis points."""
    fig, ax = plt.subplots(figsize=(5.5, 3.5))
    p = np.linspace(20, 180, 300)
    flow = np.piecewise(
        p,
        [p < 60, (p >= 60) & (p <= 140), p > 140],
        [lambda x: 1.0 * (x / 60),
         lambda x: 1.0 + 0.0 * x,
         lambda x: 1.0 + 0.005 * (x - 140)],
    )
    ax.plot(p, flow, color="black", lw=1.6)
    # Baseline operating point (blue circle)
    ax.plot(100, 1.0, "o", color="#1F4E79", markersize=9)
    ax.text(101, 1.07, "baseline", color="#1F4E79", fontsize=8)
    # Candidate points per description
    pts = {"A": (100, 1.25), "B": (95, 1.05), "C": (100, 0.6), "D": (55, 0.95)}
    for k, (x, y) in pts.items():
        ax.plot(x, y, "o", color="#C42B2B", markersize=6)
        ax.annotate(k, (x, y), xytext=(x + 3, y + 0.05), fontsize=10, fontweight="bold")
    ax.set_xlabel("Coronary perfusion pressure (mm Hg)")
    ax.set_ylabel("Coronary blood flow (relative)")
    ax.set_xlim(20, 180); ax.set_ylim(0, 1.5)
    ax.set_title("Coronary autoregulation", fontsize=10)
    _setup_axes(ax)
    _save(fig, path)


def fig_q38(path):
    """Five PV loop panels (A-E), baseline solid vs post-therapy dashed. E is correct (nitroprusside): EDV down, peak P down, SV preserved."""
    fig, axes = plt.subplots(1, 5, figsize=(10, 2.7))
    panels = ["A", "B", "C", "D", "E"]
    for i, ax in enumerate(axes):
        _pv_loop_points(ax)
        ax.set_xlim(20, 170); ax.set_ylim(-5, 160)
        ax.set_title(panels[i], fontsize=10, fontweight="bold")
        ax.tick_params(labelsize=7)
        if i == 0:
            ax.set_ylabel("Pres", fontsize=8)
        else:
            ax.set_ylabel("")
        ax.set_xlabel("Vol", fontsize=8)
        # Variants
        if i == 0:    # A: bigger overall
            _draw_dashed_loop(ax, edv=140, esv=60, sp=140, dp=85)
        elif i == 1:  # B: only afterload reduction (peak down, ESV down, EDV same)
            _draw_dashed_loop(ax, edv=120, esv=35, sp=100, dp=70)
        elif i == 2:  # C: only preload reduction (EDV down, ESV same, SP same)
            _draw_dashed_loop(ax, edv=95, esv=50, sp=120, dp=78)
        elif i == 3:  # D: opposite, EDV up SP up
            _draw_dashed_loop(ax, edv=140, esv=70, sp=140, dp=88)
        else:         # E: correct nitroprusside, EDV down, SP modestly down, SV preserved
            _draw_dashed_loop(ax, edv=95, esv=30, sp=105, dp=72)
    fig.suptitle("PV loops, baseline (solid) vs post therapy (dashed)", fontsize=10)
    _save(fig, path)


# Map display number -> figure builder
FIG_BUILDERS = {
    1: fig_q01,
    2: fig_q02,
    4: fig_q04,
    5: fig_q05,
    7: fig_q07,
    8: fig_q08,
    11: fig_q11,
    12: fig_q12,
    16: fig_q16,
    21: fig_q21,
    24: fig_q24,
    34: fig_q34,
    35: fig_q35,
    36: fig_q36,
    38: fig_q38,
}


# ---------- Helpers for pptx ----------

def add_textbox(slide, x, y, w, h, text, *, font_size=12, bold=False, italic=False,
                color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                font_name=FONT_BODY, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=GRAY_LINE, line_w=0.75, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    if rounded:
        # Tighter corner radius
        s.adjustments[0] = 0.08
    s.shadow.inherit = False
    return s


def add_footer(slide):
    add_textbox(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.3),
                "USMLE Step 1 Prep  ·  CVS Physiology",
                font_size=9, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def add_flag(slide):
    """Small red flag icon, top right (USMLE-style mark question)."""
    flag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   SLIDE_W - Inches(0.55), Inches(0.35),
                                   Inches(0.28), Inches(0.18))
    flag.fill.solid()
    flag.fill.fore_color.rgb = USMLE_RED
    flag.line.fill.background()
    flag.shadow.inherit = False
    pole = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   SLIDE_W - Inches(0.28), Inches(0.32),
                                   Inches(0.03), Inches(0.35))
    pole.fill.solid()
    pole.fill.fore_color.rgb = BLACK
    pole.line.fill.background()
    pole.shadow.inherit = False


# ---------- Slide builders ----------

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Subtle horizontal band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0),
                                   SLIDE_W, Inches(1.8))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0xF6, 0xF9, 0xFC)
    band.line.fill.background()
    band.shadow.inherit = False

    add_textbox(slide, Inches(0.7), Inches(3.05), SLIDE_W - Inches(1.4), Inches(0.9),
                "Cardiovascular Physiology Review",
                font_size=44, bold=True, color=USMLE_BLUE, font_name=FONT_HEAD,
                align=PP_ALIGN.LEFT, line_spacing=1.0)
    add_textbox(slide, Inches(0.7), Inches(3.95), SLIDE_W - Inches(1.4), Inches(0.6),
                "40 USMLE style vignettes  ·  Days 4 to 6 wrap up",
                font_size=20, color=GRAY_TEXT, align=PP_ALIGN.LEFT, line_spacing=1.0)
    add_textbox(slide, Inches(0.7), Inches(4.55), SLIDE_W - Inches(1.4), Inches(0.5),
                "Allan Bakesiga, MD  ·  USMLE Step 1 Cohort",
                font_size=14, color=GRAY_MID, italic=True, align=PP_ALIGN.LEFT)
    # Thin accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.95),
                                   Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = USMLE_RED
    line.line.fill.background()
    line.shadow.inherit = False


def build_divider_slide(prs, title="Answer Key"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(1.2),
                title, font_size=52, bold=True, color=USMLE_BLUE,
                align=PP_ALIGN.CENTER, font_name=FONT_HEAD, line_spacing=1.0)
    add_textbox(slide, Inches(0), Inches(4.2), SLIDE_W, Inches(0.6),
                "Answers 1 to 40 with explanations",
                font_size=18, color=GRAY_MID, italic=True, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.07), Inches(4.05),
                                   Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = USMLE_RED
    line.line.fill.background()
    line.shadow.inherit = False


def _option_block(slide, x, y, w, options, font_size=12):
    """Render options as a vertical list of thin-bordered rows with a light gray radio."""
    row_h = Inches(0.40)
    gap = Inches(0.06)
    for opt in options:
        # Container
        rect = add_rect(slide, x, y, w, row_h, fill=WHITE, line=GRAY_LINE,
                        line_w=0.5, rounded=True)
        # Radio circle
        radio_d = Inches(0.18)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          x + Inches(0.12), y + (row_h - radio_d) / 2,
                                          radio_d, radio_d)
        circle.fill.solid()
        circle.fill.fore_color.rgb = GRAY_LIGHT
        circle.line.color.rgb = GRAY_LINE
        circle.line.width = Pt(0.5)
        circle.shadow.inherit = False
        # Text
        add_textbox(slide, x + Inches(0.38), y, w - Inches(0.45), row_h,
                    opt, font_size=font_size, color=GRAY_TEXT,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        y = y + row_h + gap
    return y


def build_question_slide(prs, qnum, rec, fig_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Header
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(4), Inches(0.5),
                f"Question {qnum}", font_size=22, bold=True, color=USMLE_BLUE,
                font_name=FONT_HEAD, line_spacing=1.0)
    # Thin separator under header
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.82),
                                  SLIDE_W - Inches(1.0), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = GRAY_LINE
    sep.line.fill.background()
    sep.shadow.inherit = False
    add_flag(slide)

    has_fig = fig_path is not None and Path(fig_path).exists()
    stem_text = rec["stem_paraphrase"]
    options = rec["options_paraphrase"]
    stem_len = len(stem_text)

    if has_fig:
        # Two column: stem on left, figure + options on right
        left_x = Inches(0.5)
        left_w = Inches(6.5)
        right_x = Inches(7.3)
        right_w = Inches(5.55)

        # Stem text
        stem_font = 14 if stem_len < 500 else (13 if stem_len < 750 else 12)
        add_textbox(slide, left_x, Inches(1.05), left_w, Inches(5.6),
                    stem_text, font_size=stem_font, color=BLACK, line_spacing=1.30)

        # Figure
        fig_h = Inches(3.2)
        slide.shapes.add_picture(str(fig_path), right_x, Inches(1.05),
                                  width=right_w, height=fig_h)
        # Figure border (no fill so the picture shows through)
        add_rect(slide, right_x, Inches(1.05), right_w, fig_h,
                 fill=None, line=GRAY_LINE, line_w=0.75)

        # Options below figure
        opt_y = Inches(4.4)
        opt_font = 12 if len(options) <= 6 else 11
        _option_block(slide, right_x, opt_y, right_w, options, font_size=opt_font)
    else:
        # Single column: stem on top, options below
        stem_font = 15 if stem_len < 500 else (14 if stem_len < 750 else 13)
        stem_h = Inches(3.3 if stem_len < 600 else 3.8)
        add_textbox(slide, Inches(0.5), Inches(1.05), SLIDE_W - Inches(1.0), stem_h,
                    stem_text, font_size=stem_font, color=BLACK, line_spacing=1.30)
        # Options two-column if many
        opt_y = Inches(1.1) + stem_h
        if opt_y > Inches(4.6):
            opt_y = Inches(4.6)
        n = len(options)
        opt_font = 13 if n <= 5 else 12
        # Single column full width
        _option_block(slide, Inches(0.6), opt_y, Inches(10), options, font_size=opt_font)

    add_footer(slide)


def _set_no_fill(shape):
    """Set a shape's fill to noFill via raw XML."""
    sppr = shape.fill._xPr  # spPr
    # remove existing fill child
    for tag in ("a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:noFill"):
        for el in sppr.findall(qn(tag)):
            sppr.remove(el)
    nofill = etree.SubElement(sppr, qn("a:noFill"))


def _short_vignette(stem):
    """Make a one-sentence vignette brief from stem (no em dashes)."""
    # Take first sentence or first 160 chars
    s = stem.strip()
    # Find first period followed by space
    for end in (". ", "? "):
        idx = s.find(end)
        if 30 < idx < 220:
            return s[: idx + 1]
    if len(s) > 200:
        return s[:200].rsplit(" ", 1)[0] + "..."
    return s


def build_answer_slide(prs, qnum, rec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Header
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(4), Inches(0.5),
                f"Answer {qnum}", font_size=22, bold=True, color=USMLE_BLUE,
                font_name=FONT_HEAD, line_spacing=1.0)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.82),
                                  SLIDE_W - Inches(1.0), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = GRAY_LINE
    sep.line.fill.background()
    sep.shadow.inherit = False

    # Vignette line
    add_textbox(slide, Inches(0.5), Inches(0.95), SLIDE_W - Inches(1.0), Inches(0.45),
                f"Vignette: {_short_vignette(rec['stem_paraphrase'])}",
                font_size=11, italic=True, color=GRAY_MID, line_spacing=1.2)

    # Correct answer + green check
    check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(1.45),
                                    Inches(0.32), Inches(0.32))
    check.fill.solid()
    check.fill.fore_color.rgb = USMLE_GREEN
    check.line.fill.background()
    check.shadow.inherit = False
    add_textbox(slide, Inches(0.5), Inches(1.43), Inches(0.32), Inches(0.36),
                "✓", font_size=16, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.92), Inches(1.42), Inches(7), Inches(0.4),
                f"Correct answer: {rec['correct_answer']}",
                font_size=16, bold=True, color=USMLE_GREEN, line_spacing=1.1)

    # Explanation block on left
    exp = rec["explanation_gist"]
    exp_font = 12 if len(exp) < 900 else (11 if len(exp) < 1200 else 10)
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(7.4), Inches(5.0),
                exp, font_size=exp_font, color=BLACK, line_spacing=1.32)

    # Distractor box on right
    dx = Inches(8.1)
    dw = SLIDE_W - Inches(8.6)
    dy = Inches(2.0)
    dh = Inches(4.9)
    box = add_rect(slide, dx, dy, dw, dh, fill=GRAY_LIGHT, line=GRAY_LINE, line_w=0.5, rounded=True)
    # Title for the box
    add_textbox(slide, dx + Inches(0.15), dy + Inches(0.1), dw - Inches(0.3), Inches(0.4),
                "Why other choices are wrong", font_size=12, bold=True,
                color=USMLE_BLUE, font_name=FONT_HEAD)
    # Distractor entries
    notes = rec.get("distractor_notes", {}) or {}
    entries = sorted(notes.items())
    # Pick font size based on total length
    total = sum(len(k) + len(v) for k, v in entries)
    note_font = 10 if total < 600 else (9 if total < 900 else 8)

    tb = slide.shapes.add_textbox(dx + Inches(0.15), dy + Inches(0.55),
                                    dw - Inches(0.3), dh - Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    first = True
    for letter, note in entries:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = 1.22
        p.space_before = Pt(2)
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = f"{letter}. "
        r1.font.bold = True
        r1.font.size = Pt(note_font)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = USMLE_BLUE
        r2 = p.add_run()
        r2.text = note
        r2.font.bold = False
        r2.font.size = Pt(note_font)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = GRAY_TEXT

    add_footer(slide)


# ---------- Sanity: strip em dashes ----------

def scrub_dashes(rec):
    """Replace any em or en dashes in source text with commas (per Allan's rules)."""
    repl = {"—": ",", "–": ",", "—": ",", "–": ","}
    def fix(s):
        if not isinstance(s, str):
            return s
        out = s
        for k, v in repl.items():
            out = out.replace(k, v)
        return out

    for k, v in list(rec.items()):
        if isinstance(v, str):
            rec[k] = fix(v)
        elif isinstance(v, list):
            rec[k] = [fix(x) if isinstance(x, str) else x for x in v]
        elif isinstance(v, dict):
            rec[k] = {kk: fix(vv) for kk, vv in v.items()}
    return rec


# ---------- Build ----------

def main():
    with JSONL.open() as f:
        records = [json.loads(l) for l in f if l.strip()]
    records.sort(key=lambda r: r["q_num"])
    selected = [r for r in records if r["q_num"] >= 11]
    assert len(selected) == 40, f"Expected 40, got {len(selected)}"
    selected = [scrub_dashes(r) for r in selected]

    # Generate figures
    figs_generated = 0
    needed = 0
    fig_paths = {}
    for i, rec in enumerate(selected, 1):
        if rec.get("image_present"):
            needed += 1
            fig_path = FIG_DIR / f"q{i:02d}.png"
            builder = FIG_BUILDERS.get(i)
            if builder:
                try:
                    builder(fig_path)
                    figs_generated += 1
                    fig_paths[i] = fig_path
                    print(f"  fig Q{i:02d}: ok")
                except Exception as e:
                    print(f"  fig Q{i:02d}: FAILED {e}")
            else:
                print(f"  fig Q{i:02d}: no builder")

    # Build pptx
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    for i, rec in enumerate(selected, 1):
        build_question_slide(prs, i, rec, fig_path=fig_paths.get(i))
    build_divider_slide(prs, "Answer Key")
    for i, rec in enumerate(selected, 1):
        build_answer_slide(prs, i, rec)

    prs.save(str(OUT_PPTX))
    size_mb = OUT_PPTX.stat().st_size / (1024 * 1024)
    print(f"\nSaved: {OUT_PPTX}")
    print(f"Slides: {len(prs.slides)}")
    print(f"File size: {size_mb:.2f} MB")
    print(f"Figures generated: {figs_generated} of {needed} needed")


if __name__ == "__main__":
    main()
